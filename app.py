from gevent import monkey; monkey.patch_all()
import re
import io
import os
import json
import tempfile
import subprocess
import hmac
import hashlib
import threading
from uuid import uuid4
from datetime import datetime, timedelta
from dotenv import load_dotenv

load_dotenv(os.path.join(os.path.dirname(__file__), '.env'))

import anthropic
import requests as http_requests
from functools import wraps
from werkzeug.security import generate_password_hash, check_password_hash
from flask import Flask, render_template, request, jsonify, send_file, session, redirect
import db as user_db
# youtube_transcript_api כגיבוי בלבד — גישה ראשית דרך Invidious
try:
    from youtube_transcript_api import YouTubeTranscriptApi
    _YT_API_AVAILABLE = True
except ImportError:
    _YT_API_AVAILABLE = False
from deep_translator import GoogleTranslator
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)
app.secret_key = os.getenv('SECRET_KEY', 'shulchan4-secret-2026')

from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
limiter = Limiter(app, key_func=get_remote_address, default_limits=["200/minute"])

_transcript_cache = {}
_transcript_cache_lock = threading.Lock()

def _get_cached_transcript(vid):
    with _transcript_cache_lock:
        entry = _transcript_cache.get(vid)
        if entry and datetime.now() - entry['ts'] < timedelta(hours=24):
            return entry['snippets'], entry['lang']
    return None, None

def _set_cached_transcript(vid, snippets, lang):
    with _transcript_cache_lock:
        _transcript_cache[vid] = {'snippets': snippets, 'lang': lang, 'ts': datetime.now()}

COOKIES_FILE = os.path.join(os.path.dirname(__file__), 'yt_cookies.txt')
RESEND_KEY   = os.getenv('RESEND_API_KEY', '')

LS_API_KEY                = os.getenv('LS_API_KEY', '')
LS_STORE_ID               = os.getenv('LS_STORE_ID', '')
LS_VARIANT_ID             = os.getenv('LS_VARIANT_ID', '')
LS_VARIANT_UNLIMITED_ID   = os.getenv('LS_VARIANT_UNLIMITED_ID', '')
LS_WEBHOOK_SECRET         = os.getenv('LS_WEBHOOK_SECRET', '')
APP_URL                   = os.getenv('APP_URL', 'http://localhost:5002')
SUPADATA_API_KEY          = os.getenv('SUPADATA_API_KEY', '')
GUMROAD_PRO_URL           = os.getenv('GUMROAD_PRO_URL', 'https://jakovdob.gumroad.com/l/fkxlic')
GUMROAD_UNLIMITED_URL     = os.getenv('GUMROAD_UNLIMITED_URL', 'https://jakovdob.gumroad.com/l/stbsdj')
GUMROAD_PRO_PERMALINK     = 'fkxlic'
GUMROAD_UNLIMITED_PERMALINK = 'stbsdj'
FREE_LIMIT                = 3
PRO_LIMIT                 = 20
ADMIN_EMAILS      = ['jakovdob@gmail.com']
TELEGRAM_TOKEN    = os.getenv('TELEGRAM_TOKEN', '')
ADMIN_CHAT_ID     = os.getenv('ADMIN_CHAT_ID', '')

user_db.init()


def notify_ceo(text: str):
    """שולח הודעת Telegram למנכ"ל — fire and forget."""
    if not TELEGRAM_TOKEN or not ADMIN_CHAT_ID:
        return
    try:
        http_requests.post(
            f'https://api.telegram.org/bot{TELEGRAM_TOKEN}/sendMessage',
            json={'chat_id': ADMIN_CHAT_ID, 'text': text, 'parse_mode': 'HTML'},
            timeout=5,
        )
    except Exception:
        pass

def current_user():
    uid = session.get('user_id')
    if not uid:
        return None
    return user_db.get_by_id(uid)

def is_admin(u) -> bool:
    return u and u.get('email', '') in ADMIN_EMAILS


def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get('user_id'):
            return redirect('/login')
        return f(*args, **kwargs)
    return decorated


BETA_MODE = False

def pro_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        u = current_user()
        if not u:
            return jsonify({'error': 'לא מחובר', 'upgrade': False}), 401
        if BETA_MODE:
            return f(*args, **kwargs)
        if u['plan'] not in ('pro', 'unlimited') and not is_admin(u):
            return jsonify({'error': 'פיצ\'ר זה זמין לחברי Pro בלבד', 'upgrade': True}), 403
        return f(*args, **kwargs)
    return decorated


@app.route('/login', methods=['GET', 'POST'])
@limiter.limit("10/minute", methods=["POST"])
def login():
    if session.get('user_id'):
        return redirect('/')
    error = ''
    if request.method == 'POST':
        email = request.form.get('email', '').strip().lower()
        pw    = request.form.get('password', '')
        u = user_db.get_by_email(email)
        if u and check_password_hash(u['password_hash'], pw):
            session['user_id'] = u['id']
            return redirect('/')
        error = 'אימייל או סיסמה שגויים'
    return render_template('login.html', error=error)


@app.route('/register', methods=['GET', 'POST'])
@limiter.limit("5/minute", methods=["POST"])
def register():
    if session.get('user_id'):
        return redirect('/')
    error = ''
    if request.method == 'POST':
        name  = request.form.get('name', '').strip()
        email = request.form.get('email', '').strip().lower()
        pw    = request.form.get('password', '')
        if not name or not email or not pw:
            error = 'נא למלא את כל השדות'
        elif len(pw) < 6:
            error = 'הסיסמה חייבת להכיל לפחות 6 תווים'
        elif user_db.get_by_email(email):
            error = 'האימייל כבר רשום במערכת'
        else:
            try:
                uid = user_db.create(email, generate_password_hash(pw), name)
                session['user_id'] = uid
                notify_ceo(f'🆕 <b>נרשם משתמש חדש!</b>\n👤 {name}\n📧 {email}')
                return redirect('/')
            except Exception:
                error = 'שגיאה בהרשמה, נסה שוב'
    return render_template('register.html', error=error)


@app.route('/logout')
def logout():
    session.clear()
    return redirect('/login')


@app.route('/forgot-password', methods=['GET', 'POST'])
def forgot_password():
    msg = ''
    if request.method == 'POST':
        email = request.form.get('email', '').strip().lower()
        u = user_db.get_by_email(email)
        if u and RESEND_KEY:
            token = user_db.create_reset_token(u['id'])
            reset_url = f"{request.host_url.rstrip('/')}/reset-password/{token}"
            http_requests.post(
                'https://api.resend.com/emails',
                headers={'Authorization': f'Bearer {RESEND_KEY}', 'Content-Type': 'application/json'},
                json={
                    'from': 'PodSnap <onboarding@resend.dev>',
                    'to': [email],
                    'subject': 'איפוס סיסמה — PodSnap',
                    'html': f'''<div dir="rtl" style="font-family:Arial,sans-serif;max-width:480px;margin:auto">
                        <h2 style="color:#7C3AED">PodSnap — איפוס סיסמה</h2>
                        <p>שלום, קיבלנו בקשה לאיפוס הסיסמה שלך.</p>
                        <p>לחץ על הכפתור כדי לאפס (תוקף: שעה אחת):</p>
                        <a href="{reset_url}" style="display:inline-block;margin:20px 0;padding:14px 28px;background:#7C3AED;color:#fff;border-radius:8px;text-decoration:none;font-weight:700">
                            אפס סיסמה ←
                        </a>
                        <p style="color:#888;font-size:0.85rem">אם לא ביקשת איפוס, התעלם ממייל זה.</p>
                    </div>'''
                },
                timeout=10
            )
        msg = 'אם האימייל רשום במערכת — ישלח אליך קישור לאיפוס'
    return render_template('forgot_password.html', msg=msg)


@app.route('/reset-password/<token>', methods=['GET', 'POST'])
def reset_password_page(token):
    row = user_db.get_reset_token(token)
    if not row:
        return render_template('reset_password.html', error='הקישור לא תקין או שפג תוקפו', expired=True, token=token)
    error = ''
    if request.method == 'POST':
        pw  = request.form.get('password', '')
        pw2 = request.form.get('password2', '')
        if len(pw) < 6:
            error = 'הסיסמה חייבת להכיל לפחות 6 תווים'
        elif pw != pw2:
            error = 'הסיסמאות לא תואמות'
        else:
            user_db.update_password(row['user_id'], generate_password_hash(pw))
            user_db.delete_reset_token(token)
            return redirect('/login?reset=1')
    return render_template('reset_password.html', error=error, expired=False, token=token)


@app.route('/pricing')
def pricing():
    u = current_user()
    upgraded = request.args.get('upgraded') == '1'
    return render_template('pricing.html', user=u, upgraded=upgraded)


@app.route('/terms')
def terms():
    return render_template('terms.html')


@app.route('/privacy')
def privacy():
    return render_template('privacy.html')


@app.route('/create-checkout-session', methods=['POST'])
@login_required
def create_checkout():
    u = current_user()
    req_data = request.json or {}
    plan = req_data.get('plan', 'pro')

    if plan == 'unlimited':
        if u['plan'] == 'unlimited':
            return jsonify({'error': 'כבר מנוי Unlimited'}), 400
        url = GUMROAD_UNLIMITED_URL
    else:
        if u['plan'] in ('pro', 'unlimited'):
            return jsonify({'error': 'כבר מנוי Pro'}), 400
        url = GUMROAD_PRO_URL

    return jsonify({'url': url})


@app.route('/ls-webhook', methods=['POST'])
def ls_webhook():
    payload = request.get_data()
    signature = request.headers.get('X-Signature', '')
    if LS_WEBHOOK_SECRET:
        expected = hmac.new(LS_WEBHOOK_SECRET.encode(), payload, hashlib.sha256).hexdigest()
        if not hmac.compare_digest(expected, signature):
            return '', 400

    data = request.json or {}
    event = data.get('meta', {}).get('event_name', '')
    attrs = data.get('data', {}).get('attributes', {})

    if event == 'subscription_created':
        custom = data.get('meta', {}).get('custom_data', {})
        uid = int(custom.get('user_id', 0))
        if uid:
            customer_id = str(attrs.get('customer_id', ''))
            sub_id = str(data.get('data', {}).get('id', ''))
            portal_url = attrs.get('urls', {}).get('customer_portal', '')
            webhook_variant = str(attrs.get('variant_id', ''))
            is_unlimited = LS_VARIANT_UNLIMITED_ID and webhook_variant == str(LS_VARIANT_UNLIMITED_ID)
            new_plan = 'unlimited' if is_unlimited else 'pro'
            user_db.set_pro(uid, customer_id, sub_id, portal_url)
            if is_unlimited:
                user_db.set_plan(uid, 'unlimited')
            u = user_db.get_by_id(uid)
            name = u.get('name', '') if u else ''
            email = u.get('email', '') if u else ''
            price_str = '₪79/חודש' if is_unlimited else '₪29/חודש'
            notify_ceo(f'💰 <b>תשלום חדש! {new_plan.capitalize()}</b>\n👤 {name}\n📧 {email}\n💳 {price_str}')

    elif event in ('subscription_cancelled', 'subscription_expired', 'subscription_paused'):
        customer_id = str(attrs.get('customer_id', ''))
        u = user_db.get_by_payment_customer(customer_id)
        if u:
            user_db.set_free(u['id'])
            notify_ceo(f'😞 <b>ביטול מנוי</b>\n👤 {u.get("name","")}\n📧 {u.get("email","")}\nאירוע: {event}')

    return '', 200


@app.route('/gumroad-webhook', methods=['POST'])
def gumroad_webhook():
    permalink = request.form.get('product_permalink', '')
    email = request.form.get('email', '')
    refunded = request.form.get('refunded', 'false').lower() == 'true'
    if not email:
        return 'ok', 200
    u = user_db.get_by_email(email)
    if not u:
        return 'ok', 200
    if refunded:
        user_db.set_plan(u['id'], 'free')
        notify_ceo(f'😞 <b>ביטול/החזר Gumroad</b>\n📧 {email}')
    elif permalink == GUMROAD_UNLIMITED_PERMALINK:
        user_db.set_plan(u['id'], 'unlimited')
        notify_ceo(f'💰 <b>Gumroad Unlimited!</b>\n📧 {email}')
    elif permalink == GUMROAD_PRO_PERMALINK:
        user_db.set_plan(u['id'], 'pro')
        notify_ceo(f'💰 <b>Gumroad Pro!</b>\n📧 {email}')
    return 'ok', 200


@app.route('/billing-portal', methods=['POST'])
@login_required
def billing_portal():
    u = current_user()
    portal_url = u['portal_url'] if u['portal_url'] else None
    if not portal_url:
        return jsonify({'error': 'קישור לניהול מנוי לא זמין, פנה אלינו'}), 400
    return jsonify({'url': portal_url})


# ── yt-dlp availability ────────────────────────────────────
try:
    import yt_dlp
    _YTDLP_AVAILABLE = True
except ImportError:
    _YTDLP_AVAILABLE = False

# ── Invidious instances (ניסיון בסדר הזה) ─────────────────
INVIDIOUS_INSTANCES = [
    'https://inv.nadeko.net',
    'https://invidious.privacydev.net',
    'https://invidious.jing.rocks',
    'https://yewtu.be',
    'https://invidious.fdn.fr',
    'https://yt.cdaut.de',
    'https://invidious.drgns.space',
    'https://inv.zzls.xyz',
    'https://invidious.privacyredirect.com',
    'https://invidious.perennialte.ch',
]

def fetch_via_supadata(video_id):
    """Supadata API — שירות חיצוני שעובד מכל שרת"""
    if not SUPADATA_API_KEY:
        return None, None

    class Snippet:
        def __init__(self, start, text):
            self.start = start
            self.text  = text

    try:
        r = http_requests.get(
            'https://api.supadata.ai/v1/youtube/transcript',
            params={'url': f'https://www.youtube.com/watch?v={video_id}', 'text': 'false'},
            headers={'x-api-key': SUPADATA_API_KEY},
            timeout=30
        )
        print(f'[supadata] status={r.status_code}')
        if r.status_code != 200:
            print(f'[supadata] error: {r.text[:200]}')
            return None, None

        data = r.json()
        content = data.get('content', [])
        lang = data.get('lang', 'he')

        if not content:
            print('[supadata] empty content')
            return None, None

        snippets = [Snippet(item.get('offset', 0) / 1000, item.get('text', ''))
                    for item in content if item.get('text', '').strip()]
        print(f'[supadata] {len(snippets)} snippets ({lang})')
        return snippets, lang

    except Exception as e:
        print(f'[supadata] error: {e}')
        return None, None


def fetch_via_innertube(video_id):
    """YouTube InnerTube API (Android client) — פחות חסום מ-cloud servers"""
    class Snippet:
        def __init__(self, start, text):
            self.start = start
            self.text  = text

    CLIENTS = [
        {
            'name': 'ANDROID',
            'headers': {
                'X-YouTube-Client-Name': '3',
                'X-YouTube-Client-Version': '17.10.35',
                'User-Agent': 'com.google.android.youtube/17.10.35 (Linux; U; Android 11) gzip',
                'Content-Type': 'application/json',
            },
            'payload': {
                'context': {'client': {'clientName': 'ANDROID', 'clientVersion': '17.10.35', 'androidSdkVersion': 30, 'hl': 'en', 'gl': 'US'}},
            },
            'key': 'AIzaSyA8eiZmM1FaDVjRy-df2KTyQ_vz_yYM394',
        },
        {
            'name': 'TVHTML5',
            'headers': {
                'X-YouTube-Client-Name': '7',
                'X-YouTube-Client-Version': '7.20201028',
                'User-Agent': 'Mozilla/5.0 (SMART-TV; LINUX; Tizen 5.0) AppleWebKit/538.1',
                'Content-Type': 'application/json',
            },
            'payload': {
                'context': {'client': {'clientName': 'TVHTML5', 'clientVersion': '7.20201028', 'hl': 'en', 'gl': 'US'}},
            },
            'key': 'AIzaSyAO_FJ2SlqU8Q4STEHLGCilw_Y9_11qcW8',
        },
    ]

    for client in CLIENTS:
        try:
            payload = dict(client['payload'])
            payload['videoId'] = video_id
            r = http_requests.post(
                f"https://www.youtube.com/youtubei/v1/player?key={client['key']}",
                json=payload, headers=client['headers'], timeout=15
            )
            if r.status_code != 200:
                print(f"[innertube/{client['name']}] status {r.status_code}")
                continue

            data = r.json()
            captions = (data.get('captions', {})
                           .get('playerCaptionsTracklistRenderer', {})
                           .get('captionTracks', []))

            if not captions:
                print(f"[innertube/{client['name']}] no captionTracks")
                continue

            chosen = None
            for lang in ['he', 'iw', 'en']:
                for cap in captions:
                    if lang in cap.get('languageCode', '').lower():
                        chosen = cap
                        break
                if chosen:
                    break
            if not chosen:
                chosen = captions[0]

            base_url = chosen.get('baseUrl', '')
            if not base_url:
                continue

            sep = '&' if '?' in base_url else '?'
            for fmt in ['json3', 'vtt']:
                tr = http_requests.get(base_url + sep + f'fmt={fmt}', timeout=10)
                if tr.status_code != 200:
                    continue
                if fmt == 'json3':
                    events = tr.json().get('events', [])
                    snippets = []
                    for ev in events:
                        if 'segs' not in ev:
                            continue
                        s = ev.get('tStartMs', 0) / 1000
                        txt = ' '.join(sg.get('utf8', '') for sg in ev['segs']).strip()
                        if txt:
                            snippets.append(Snippet(s, txt))
                else:
                    snippets = _parse_vtt(tr.text)

                if snippets:
                    lang_code = chosen.get('languageCode', 'he')
                    print(f"[innertube/{client['name']}] {len(snippets)} snippets ({lang_code})")
                    return snippets, lang_code

        except Exception as e:
            print(f"[innertube/{client['name']}] error: {e}")

    return None, None


def fetch_via_page_parse(video_id):
    """מחלץ תמלול מדף YouTube — עובד גם מ-cloud עם CONSENT cookie"""
    HEADERS = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36',
        'Accept-Language': 'en-US,en;q=0.9',
        'Cookie': 'SOCS=CAI; CONSENT=YES+1',
    }

    class Snippet:
        def __init__(self, start, text):
            self.start = start
            self.text  = text

    try:
        r = http_requests.get(f'https://www.youtube.com/watch?v={video_id}', headers=HEADERS, timeout=15)
        if r.status_code != 200:
            return None, None

        page = r.text
        print(f'[page-parse] status={r.status_code} size={len(page)} consent={"SOCS" in page or "consent" in page.lower()}')
        idx = page.find('"captionTracks"')
        if idx == -1:
            print('[page-parse] no captionTracks in page')
            return None, None

        start = page.find('[', idx)
        depth, end = 0, start
        for i, ch in enumerate(page[start:], start):
            if ch == '[': depth += 1
            elif ch == ']':
                depth -= 1
                if depth == 0:
                    end = i
                    break

        tracks = json.loads(page[start:end + 1])
        if not tracks:
            return None, None

        chosen = None
        for lang in ['he', 'iw', 'en']:
            for t in tracks:
                if lang in t.get('languageCode', '').lower():
                    chosen = t
                    break
            if chosen:
                break
        if not chosen:
            chosen = tracks[0]

        base_url = chosen.get('baseUrl', '')
        if not base_url:
            return None, None

        sep = '&' if '?' in base_url else '?'
        for fmt in ['vtt', 'json3']:
            tr = http_requests.get(base_url + sep + f'fmt={fmt}', headers=HEADERS, timeout=10)
            if tr.status_code != 200:
                continue
            if fmt == 'vtt':
                snippets = _parse_vtt(tr.text)
                if snippets:
                    lang_code = chosen.get('languageCode', 'he')
                    print(f'[page-parse vtt] {len(snippets)} snippets ({lang_code})')
                    return snippets, lang_code
            else:
                events = tr.json().get('events', [])
                snippets = []
                for ev in events:
                    if 'segs' not in ev:
                        continue
                    s = ev.get('tStartMs', 0) / 1000
                    txt = ' '.join(sg.get('utf8', '') for sg in ev['segs']).strip()
                    if txt:
                        snippets.append(Snippet(s, txt))
                if snippets:
                    lang_code = chosen.get('languageCode', 'he')
                    print(f'[page-parse json3] {len(snippets)} snippets ({lang_code})')
                    return snippets, lang_code
    except Exception as e:
        print(f'[page-parse] error: {e}')

    return None, None


def fetch_via_timedtext(video_id):
    """YouTube timedtext API ישיר — עובד לסרטונים ציבוריים עם כתוביות"""
    HEADERS = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36',
        'Accept-Language': 'he-IL,he;q=0.9,en;q=0.8',
        'Referer': 'https://www.youtube.com/',
    }
    for lang in ['iw', 'he', 'en', 'en-US']:
        try:
            url = f'https://www.youtube.com/api/timedtext?v={video_id}&lang={lang}&fmt=vtt'
            r = http_requests.get(url, headers=HEADERS, timeout=10)
            if r.status_code == 200 and len(r.text.strip()) > 100:
                snippets = _parse_vtt(r.text)
                if snippets:
                    print(f"[timedtext] {len(snippets)} snippets ({lang})")
                    return snippets, lang
        except Exception as e:
            print(f"[timedtext] {lang}: {e}")
    return None, None

def _parse_vtt(vtt_text):
    """ממיר WebVTT לרשימת snippets עם .start ו-.text"""
    class Snippet:
        def __init__(self, start, text):
            self.start = start
            self.text  = text

    snippets = []
    lines = vtt_text.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        # זיהוי שורת זמן: 00:00:00.000 --> 00:00:03.000
        if '-->' in line:
            time_part = line.split('-->')[0].strip()
            parts = time_part.split(':')
            try:
                if len(parts) == 3:
                    h, m, s = parts
                    start_sec = int(h)*3600 + int(m)*60 + float(s.replace(',', '.'))
                else:
                    m, s = parts
                    start_sec = int(m)*60 + float(s.replace(',', '.'))
            except:
                i += 1
                continue
            # אסוף את שורות הטקסט
            i += 1
            text_lines = []
            while i < len(lines) and lines[i].strip():
                t = lines[i].strip()
                # נקה תגיות HTML כמו <c>, </c>
                import re as _re
                t = _re.sub(r'<[^>]+>', '', t)
                if t:
                    text_lines.append(t)
                i += 1
            if text_lines:
                text = ' '.join(text_lines)
                # הסר כפילויות של הפסקה האחרונה
                if snippets and snippets[-1].text == text:
                    i += 1
                    continue
                snippets.append(Snippet(start_sec, text))
        i += 1
    return snippets

def fetch_via_ytdlp(video_id):
    """מוריד תמלול דרך yt-dlp — הכי אמין, תומך בכתוביות אוטומטיות"""
    if not _YTDLP_AVAILABLE:
        return None, None

    class Snippet:
        def __init__(self, start, text):
            self.start = start
            self.text  = text

    ydl_opts = {
        'skip_download': True,
        'writesubtitles': True,
        'writeautomaticsub': True,
        'subtitleslangs': ['he', 'iw', 'en'],
        'subtitlesformat': 'vtt',
        'quiet': True,
        'no_warnings': True,
    }
    if os.path.exists(COOKIES_FILE):
        ydl_opts['cookiefile'] = COOKIES_FILE

    url = f'https://www.youtube.com/watch?v={video_id}'
    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)
            subs = info.get('subtitles', {})
            auto_subs = info.get('automatic_captions', {})

            chosen_lang = None
            chosen_data = None
            for lang in ['he', 'iw', 'en']:
                if lang in subs:
                    chosen_lang = lang
                    chosen_data = subs[lang]
                    break
            if not chosen_data:
                for lang in ['he', 'iw', 'en']:
                    if lang in auto_subs:
                        chosen_lang = lang
                        chosen_data = auto_subs[lang]
                        break
            if not chosen_data and subs:
                chosen_lang = list(subs.keys())[0]
                chosen_data = subs[chosen_lang]
            if not chosen_data and auto_subs:
                chosen_lang = list(auto_subs.keys())[0]
                chosen_data = auto_subs[chosen_lang]

            if not chosen_data:
                return None, None

            vtt_url = None
            for fmt in chosen_data:
                if fmt.get('ext') == 'vtt':
                    vtt_url = fmt.get('url')
                    break
            if not vtt_url and chosen_data:
                base_url = chosen_data[0].get('url', '')
                if 'youtube.com/api/timedtext' in base_url:
                    sep = '&' if '?' in base_url else '?'
                    vtt_url = base_url + sep + 'fmt=vtt'
                else:
                    vtt_url = base_url

            if not vtt_url:
                return None, None

            rv = http_requests.get(vtt_url, timeout=15)
            if rv.status_code != 200:
                return None, None

            snippets = _parse_vtt(rv.text)
            if snippets:
                print(f"[yt-dlp] {len(snippets)} snippets ({chosen_lang})")
                return snippets, chosen_lang

    except Exception as e:
        print(f"[yt-dlp] error: {e}")

    return None, None


def fetch_via_invidious(video_id):
    """מוריד תמלול דרך Invidious API — עובד מכל IP"""
    HEADERS_INV = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) Chrome/122.0',
        'Accept': 'application/json',
    }
    lang_priority = ['he', 'iw', 'Hebrew', 'en', 'English']

    for base in INVIDIOUS_INSTANCES:
        try:
            # שלב 1: קבל רשימת כתוביות
            r = http_requests.get(
                f'{base}/api/v1/captions/{video_id}',
                headers=HEADERS_INV, timeout=10)
            if r.status_code != 200:
                continue

            captions = r.json().get('captions', [])
            if not captions:
                continue

            # שלב 2: בחר שפה לפי עדיפות
            chosen = None
            for lang in lang_priority:
                for cap in captions:
                    if lang.lower() in cap.get('language_code','').lower() or \
                       lang.lower() in cap.get('label','').lower():
                        chosen = cap
                        break
                if chosen:
                    break
            if not chosen:
                chosen = captions[0]  # כל שפה שיש

            # שלב 3: הורד את קובץ ה-VTT
            cap_url = chosen.get('url', '')
            if not cap_url.startswith('http'):
                cap_url = base + cap_url
            # הוסף format=vtt
            if '?' in cap_url:
                cap_url += '&format=vtt'
            else:
                cap_url += '?format=vtt'

            rv = http_requests.get(cap_url, headers=HEADERS_INV, timeout=15)
            if rv.status_code != 200:
                continue

            snippets = _parse_vtt(rv.text)
            if snippets:
                lang_code = chosen.get('language_code', 'he')
                print(f"[Invidious] {base} → {len(snippets)} snippets ({lang_code})")
                return snippets, lang_code

        except Exception as e:
            print(f"[Invidious] {base} error: {e}")
            continue

    return None, None

def extract_video_id(url):
    patterns = [
        r'(?:v=)([a-zA-Z0-9_-]{11})',
        r'(?:youtu\.be/)([a-zA-Z0-9_-]{11})',
        r'(?:embed/)([a-zA-Z0-9_-]{11})',
    ]
    for p in patterns:
        m = re.search(p, url)
        if m:
            return m.group(1)
    if len(url) == 11:
        return url
    return None


def seconds_to_time(s):
    s = int(s)
    h = s // 3600
    m = (s % 3600) // 60
    sec = s % 60
    if h:
        return f"{h}:{m:02d}:{sec:02d}"
    return f"{m}:{sec:02d}"


def group_into_sections(snippets, section_minutes=5):
    """קיבוץ קטעי תמלול לפרקי זמן של X דקות"""
    sections = []
    section_seconds = section_minutes * 60
    current_section = []
    current_start = 0

    for snip in snippets:
        t = snip.start
        if t >= current_start + section_seconds and current_section:
            sections.append({
                'start': current_start,
                'start_fmt': seconds_to_time(current_start),
                'end_fmt': seconds_to_time(t),
                'text': ' '.join(current_section),
            })
            current_section = []
            current_start = t
        current_section.append(snip.text)

    if current_section:
        sections.append({
            'start': current_start,
            'start_fmt': seconds_to_time(current_start),
            'end_fmt': '',
            'text': ' '.join(current_section),
        })

    return sections


@app.route('/')
def index():
    u = current_user()
    if not u:
        return render_template('landing.html')
    usage = user_db.month_usage(u['id'])
    upgraded = request.args.get('upgraded') == '1'
    li_connected = bool(LINKEDIN_ADMIN_TOKEN) or bool(u.get('linkedin_urn'))
    return render_template('index.html', user=u, usage=usage, free_limit=FREE_LIMIT, upgraded=upgraded,
                           is_admin_user=is_admin(u), li_connected=li_connected)


def _fetch_transcript_for_video(vid):
    cached_snippets, cached_lang = _get_cached_transcript(vid)
    if cached_snippets:
        return cached_snippets, cached_lang

    snippets, lang = fetch_via_supadata(vid)
    if not snippets:
        snippets, lang = fetch_via_innertube(vid)
    if not snippets:
        snippets, lang = fetch_via_page_parse(vid)
    if not snippets:
        snippets, lang = fetch_via_timedtext(vid)
    if not snippets and _YT_API_AVAILABLE:
        try:
            api = YouTubeTranscriptApi()
            for _langs in [['iw', 'he'], ['en']]:
                try:
                    transcript = api.fetch(vid, languages=_langs)
                    snippets = list(transcript)
                    if snippets:
                        lang = _langs[0]
                        break
                except Exception:
                    continue
        except Exception as e:
            print(f"[yt-api] {e}")
            snippets = None
    if not snippets:
        snippets, lang = fetch_via_invidious(vid)
    if not snippets:
        snippets, lang = fetch_via_ytdlp(vid)

    if snippets:
        _set_cached_transcript(vid, snippets, lang)
    return snippets, lang


def _process_job(job_id, video_url, uid):
    try:
        user_db.update_job(job_id, 'processing')
        vid = extract_video_id(video_url)
        if not vid:
            user_db.update_job(job_id, 'error', {'error': 'קישור יוטיוב לא תקין'})
            return
        snippets, lang = _fetch_transcript_for_video(vid)
        if not snippets:
            user_db.update_job(job_id, 'error', {'error': f'לא ניתן להוריד תמלול — ytdlp={_YTDLP_AVAILABLE}'})
            return
        sections = group_into_sections(snippets, section_minutes=5)
        full_text = ' '.join(s['text'] for s in sections)
        result = {
            'video_id': vid,
            'sections': sections,
            'total_sections': len(sections),
            'duration': seconds_to_time(snippets[-1].start) if snippets else '0:00',
            'full_text': full_text,
            'lang': lang,
            'translated': False,
        }
        user_db.update_job(job_id, 'done', result)
    except Exception as e:
        import traceback
        err_msg = str(e)
        user_db.update_job(job_id, 'error', {'error': err_msg})
        notify_ceo(f'⚠️ Job failed\nJob: {job_id}\nUser: {uid}\nURL: {video_url}\n{err_msg}')


@app.route('/api/transcript', methods=['POST'])
@login_required
@limiter.limit("3/minute", key_func=lambda: str(session.get('user_id', get_remote_address())))
def get_transcript():
    u = current_user()
    if not is_admin(u):
        usage = user_db.month_usage(u['id'])
        if u['plan'] == 'free' and usage >= FREE_LIMIT:
            return jsonify({'error': f'הגעת ל-{FREE_LIMIT} פרקים חינמיים החודש — שדרג ל-Pro', 'upgrade': True}), 403
        if u['plan'] == 'pro' and usage >= PRO_LIMIT:
            return jsonify({'error': f'הגעת ל-{PRO_LIMIT} פרקים ב-Pro החודש — שדרג ל-Unlimited', 'upgrade': True, 'upgrade_to': 'unlimited'}), 403
    data = request.json or {}
    url = data.get('url', '').strip()
    if not url:
        return jsonify({'error': 'נא להזין קישור'}), 400
    if not extract_video_id(url):
        return jsonify({'error': 'קישור יוטיוב לא תקין'}), 400

    job_id = uuid4().hex
    user_db.create_job(job_id, u['id'], url)
    user_db.inc_usage(u['id'])
    threading.Thread(target=_process_job, args=(job_id, url, u['id']), daemon=True).start()
    return jsonify({'job_id': job_id})


@app.route('/api/job/<job_id>')
@login_required
def get_job(job_id):
    job = user_db.get_job(job_id)
    if not job:
        return jsonify({'error': 'job not found'}), 404
    if str(job['user_id']) != str(session.get('user_id')) and not is_admin(current_user()):
        return jsonify({'error': 'Unauthorized'}), 403
    return jsonify({'status': job['status'], 'result': job['result']})


@app.route('/api/translate', methods=['POST'])
@login_required
def translate_sections():
    data = request.json or {}
    sections = data.get('sections', [])
    if not sections:
        return jsonify({'error': 'אין קטעים לתרגום'}), 400

    try:
        translator = GoogleTranslator(source='auto', target='iw')
        for sec in sections:
            chunk = sec['text'][:4500]
            sec['text'] = translator.translate(chunk) or sec['text']
    except Exception as e:
        return jsonify({'error': f'שגיאת תרגום: {str(e)}'}), 502

    full_text = ' '.join(s['text'] for s in sections)
    return jsonify({'sections': sections, 'full_text': full_text})


@app.route('/api/pptx', methods=['POST'])
@login_required
def make_pptx():
    from pptx.util import Emu
    from lxml import etree

    data = request.json or {}
    sections = data.get('sections', [])
    title = data.get('title', 'סיכום פרק')
    if not sections:
        return jsonify({'error': 'אין קטעים'}), 400

    # צבעים
    NAVY    = RGBColor(11, 20, 50)
    NAVY2   = RGBColor(18, 32, 75)
    ACCENT  = RGBColor(0, 180, 216)
    RED     = RGBColor(233, 69, 96)
    WHITE   = RGBColor(255, 255, 255)
    LGRAY   = RGBColor(200, 210, 230)
    GOLD    = RGBColor(255, 193, 7)
    BOX_BG  = RGBColor(25, 40, 90)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def set_bg(slide, color=NAVY):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=0):
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(line_width)
        else:
            shape.line.fill.background()
        return shape

    def add_tb(slide, text, l, t, w, h, size, bold=False, color=WHITE,
               align=PP_ALIGN.RIGHT, italic=False, font='Calibri'):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
        return tb

    def text_to_bullets(text, max_bullets=5):
        sentences = [s.strip() for s in re.split(r'[.!?।]', text) if len(s.strip()) > 15]
        return sentences[:max_bullets]

    # ── שקף כותרת ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    # פס עליון צבעוני
    add_rect(slide, 0, 0, 13.33, 0.12, ACCENT)
    # פס תחתון
    add_rect(slide, 0, 7.38, 13.33, 0.12, RED)

    # קו קישוט אנכי שמאלי
    add_rect(slide, 0.3, 1.5, 0.06, 4.5, ACCENT)

    # כותרת
    add_tb(slide, f'"{title}"', 0.6, 1.6, 12.0, 2.0, 38, bold=True,
           color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    add_tb(slide, '— שולחן ארבע —', 0.6, 3.7, 12.0, 0.8, 20,
           color=ACCENT, align=PP_ALIGN.CENTER)

    # מידע
    add_rect(slide, 4.5, 4.7, 4.33, 0.8, BOX_BG)
    add_tb(slide, f'📌  {len(sections)} קטעים', 4.5, 4.75, 4.33, 0.7, 18,
           color=GOLD, align=PP_ALIGN.CENTER, bold=True)

    # ── שקף לכל קטע ──
    for i, sec in enumerate(sections):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, NAVY)

        # פסים עיצוביים
        add_rect(slide, 0, 0, 13.33, 0.10, ACCENT)
        add_rect(slide, 0, 7.40, 13.33, 0.10, RED)
        add_rect(slide, 0, 0.10, 13.33, 1.0, NAVY2)

        # מספר קטע + זמן
        time_range = f"{sec['start_fmt']} – {sec['end_fmt']}" if sec.get('end_fmt') else sec['start_fmt']
        add_rect(slide, 0.3, 0.22, 1.0, 0.55, ACCENT)
        add_tb(slide, f'{i+1}', 0.3, 0.22, 1.0, 0.55, 22, bold=True,
               color=NAVY, align=PP_ALIGN.CENTER)
        add_tb(slide, f'⏱  {time_range}', 1.5, 0.25, 5.0, 0.5, 14,
               color=LGRAY, align=PP_ALIGN.RIGHT)

        # תצוגת תוכן כ-bullets
        bullets = text_to_bullets(sec['text'], max_bullets=5)
        if not bullets:
            bullets = [sec['text'][:200]]

        bullet_icons = ['🔹', '🔸', '💡', '📌', '▶']
        box_h = 1.0
        for j, bullet in enumerate(bullets):
            y = 1.3 + j * box_h
            add_rect(slide, 0.3, y, 12.73, 0.85, BOX_BG, line_color=ACCENT, line_width=0.5)
            text_short = bullet[:140] + ('...' if len(bullet) > 140 else '')
            add_tb(slide, f'{bullet_icons[j % 5]}  {text_short}',
                   0.45, y + 0.05, 12.43, 0.75, 15,
                   color=WHITE, align=PP_ALIGN.RIGHT)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    safe_title = re.sub(r'[^\w\s-]', '', title)[:40].strip() or 'presentation'
    return send_file(buf, as_attachment=True,
                     download_name=f'{safe_title}.pptx',
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')


def _build_smart_pptx_bytes(full_text, title):
    """יוצר PPTX ומחזיר bytes + שם קובץ בטוח"""
    from pptx.util import Emu
    import json as _json

    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        raise ValueError('חסר API key')

    client = anthropic.Anthropic(api_key=api_key)
    prompt = f"""קרא את התמלול הבא וזהה את הנושאים המרכזיים שנדונו.
לכל נושא תן:
1. כותרת קצרה (עד 8 מילים)
2. 3-4 נקודות עיקריות (כל אחת עד 20 מילים)

החזר בפורמט JSON בלבד, ללא הסברים נוספים:
{{
  "topics": [
    {{
      "title": "כותרת הנושא",
      "points": ["נקודה 1", "נקודה 2", "נקודה 3"]
    }}
  ]
}}

זהה עד 8 נושאים מרכזיים.

התמלול:
{full_text[:12000]}"""

    msg = client.messages.create(
        model='claude-haiku-4-5-20251001',
        max_tokens=2000,
        messages=[{'role': 'user', 'content': prompt}]
    )
    raw = msg.content[0].text.strip()
    json_match = re.search(r'\{.*\}', raw, re.DOTALL)
    if not json_match:
        raise ValueError('לא נמצא JSON')
    topics_data = _json.loads(json_match.group())['topics']

    BG      = RGBColor(2,  8, 18)
    BG2     = RGBColor(4, 14, 32)
    CYAN    = RGBColor(0, 255, 255)
    CYAN2   = RGBColor(0, 180, 220)
    CYAN3   = RGBColor(0,  45,  60)
    GREEN   = RGBColor(0, 255, 120)
    WHITE   = RGBColor(220, 240, 255)
    DIM     = RGBColor(80, 140, 160)
    NEONS   = [
        RGBColor(0,255,255), RGBColor(255,0,200), RGBColor(0,255,120),
        RGBColor(255,220,0), RGBColor(100,180,255), RGBColor(255,100,100),
        RGBColor(180,255,100), RGBColor(200,100,255)
    ]

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG

    def add_rect(slide, l, t, w, h, fc, lc=None, lw=0):
        s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = fc
        if lc:
            s.line.color.rgb = lc
            s.line.width = Pt(lw)
        else:
            s.line.fill.background()
        return s

    def add_tb(slide, text, l, t, w, h, size, bold=False, color=WHITE,
               align=PP_ALIGN.RIGHT, italic=False, font='Consolas'):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
        return tb

    def draw_grid(slide, step=0.7):
        x = 0.0
        while x <= 13.33:
            add_rect(slide, x, 0, 0.01, 7.5, CYAN3)
            x += step
        y = 0.0
        while y <= 7.5:
            add_rect(slide, 0, y, 13.33, 0.01, CYAN3)
            y += step

    def draw_corners(slide, color, sz=0.45, th=0.045):
        m = 0.2
        for lx, ly, dx, dy in [
            (m, m, 1, 1), (13.33-m-sz, m, -1, 1),
            (m, 7.5-m-th, 1, -1), (13.33-m-sz, 7.5-m-th, -1, -1)
        ]:
            add_rect(slide, lx, ly if dy==1 else ly, sz, th, color)
            add_rect(slide, lx if dx==1 else lx+sz-th, m if dy==1 else 7.5-m-sz, th, sz, color)

    def glow_box(slide, l, t, w, h, color):
        c1 = RGBColor(color[0]//5, color[1]//5, color[2]//5)
        c2 = RGBColor(color[0]//2, color[1]//2, color[2]//2)
        add_rect(slide, l-0.05, t-0.05, w+0.10, h+0.10, BG, lc=c1, lw=3)
        add_rect(slide, l-0.025, t-0.025, w+0.05, h+0.05, BG, lc=c2, lw=2)
        add_rect(slide, l, t, w, h, BG2, lc=color, lw=1.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    draw_grid(slide)
    draw_corners(slide, CYAN)
    add_rect(slide, 0, 3.55, 13.33, 0.025, CYAN)
    add_tb(slide, 'BRIEFING  //  CLASSIFIED', 0.5, 0.45, 12.33, 0.55, 11,
           color=DIM, align=PP_ALIGN.CENTER)
    add_tb(slide, f'// {title} //', 0.5, 1.15, 12.33, 2.3, 30,
           bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    add_tb(slide, '━' * 52, 0.5, 3.25, 12.33, 0.4, 10,
           color=CYAN2, align=PP_ALIGN.CENTER)
    add_tb(slide, f'[ שולחן ארבע ]   TOPICS: {len(topics_data):02d}   STATUS: ACTIVE',
           0.5, 3.7, 12.33, 0.6, 14, color=GREEN, align=PP_ALIGN.CENTER)
    add_rect(slide, 0, 6.82, 13.33, 0.02, CYAN2)
    add_tb(slide, 'SYSTEM ONLINE  ◈  AUDIO PROCESSED  ◈  READY',
           0.5, 6.86, 12.33, 0.45, 10, color=DIM, align=PP_ALIGN.CENTER)

    for i, topic in enumerate(topics_data):
        neon = NEONS[i % len(NEONS)]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        draw_grid(slide, step=0.85)
        draw_corners(slide, neon)
        add_rect(slide, 0, 1.25, 13.33, 0.02, neon)
        add_rect(slide, 0, 0, 13.33, 1.22, BG2)
        add_rect(slide, 0, 0, 0.07, 1.22, neon)
        glow_box(slide, 0.2, 0.2, 0.78, 0.78, neon)
        add_tb(slide, f'{i+1:02d}', 0.2, 0.22, 0.78, 0.72, 26,
               bold=True, color=neon, align=PP_ALIGN.CENTER)
        add_tb(slide, f'SEGMENT_{i+1:02d}.DAT', 0.12, 0.06, 5.0, 0.35, 8,
               color=DIM, align=PP_ALIGN.LEFT)
        add_tb(slide, f'◈  {topic["title"]}', 1.1, 0.24, 12.0, 0.75, 22,
               bold=True, color=WHITE, align=PP_ALIGN.RIGHT, font='Calibri')

        points = topic.get('points', [])
        n = max(len(points), 1)
        box_h = min(1.18, (7.5 - 1.42) / n)
        for j, point in enumerate(points):
            y = 1.35 + j * box_h
            glow_box(slide, 0.22, y, 12.89, box_h - 0.12, neon)
            add_rect(slide, 0.22, y, 0.06, box_h - 0.12, neon)
            add_tb(slide, f'{j+1:02d}', 0.3, y+0.06, 0.55, box_h-0.2, 11,
                   color=neon, align=PP_ALIGN.CENTER)
            tc = WHITE if j % 2 == 0 else RGBColor(
                min(255, neon[0]+60), min(255, neon[1]+60), min(255, neon[2]+60))
            add_tb(slide, point, 0.95, y+0.07, 12.0, box_h-0.2, 15,
                   color=tc, align=PP_ALIGN.RIGHT, font='Calibri')

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    safe = re.sub(r'[^\w\s-]', '', title)[:40].strip() or 'summary'
    return buf.getvalue(), safe


@app.route('/api/smart_pptx', methods=['POST'])
@login_required
@pro_required
def smart_pptx():
    data = request.json or {}
    full_text = data.get('full_text', '').strip()
    title = data.get('title', 'סיכום פרק')
    if not full_text:
        return jsonify({'error': 'אין תמלול'}), 400
    try:
        pptx_bytes, safe = _build_smart_pptx_bytes(full_text, title)
    except Exception as e:
        return jsonify({'error': str(e)}), 502
    return send_file(io.BytesIO(pptx_bytes), as_attachment=True,
                     download_name=f'{safe}_topics.pptx',
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')


@app.route('/api/smart_pptx_base64', methods=['POST'])
@login_required
@pro_required
def smart_pptx_base64():
    import base64
    data = request.json or {}
    full_text = data.get('full_text', '').strip()
    title = data.get('title', 'סיכום פרק')
    if not full_text:
        return jsonify({'error': 'אין תמלול'}), 400
    try:
        pptx_bytes, safe = _build_smart_pptx_bytes(full_text, title)
    except Exception as e:
        return jsonify({'error': str(e)}), 502
    return jsonify({
        'filename': f'{safe}_topics.pptx',
        'base64': base64.b64encode(pptx_bytes).decode('utf-8')
    })


@app.route('/api/js_pptx', methods=['POST'])
@login_required
@pro_required
def js_pptx():
    data = request.json or {}
    full_text = data.get('full_text', '').strip()
    title = data.get('title', 'סיכום פרק')

    if not full_text:
        return jsonify({'error': 'אין תמלול'}), 400

    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        return jsonify({'error': 'חסר API key'}), 500

    client = anthropic.Anthropic(api_key=api_key)
    prompt = f"""קרא את התמלול וזהה עד 8 נושאים מרכזיים.
לכל נושא תן כותרת קצרה (עד 8 מילים) ו-4 נקודות (כל אחת משפט אחד עד 25 מילים).
החזר JSON בלבד:
{{"topics":[{{"title":"...","points":["...","...","...","..."]}}]}}

התמלול:
{full_text[:12000]}"""

    try:
        msg = client.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=2000,
            messages=[{'role': 'user', 'content': prompt}]
        )
        raw = msg.content[0].text.strip()
        m = re.search(r'\{.*\}', raw, re.DOTALL)
        topics_data = json.loads(m.group())['topics']
    except Exception as e:
        return jsonify({'error': f'שגיאת Claude: {str(e)}'}), 502

    # כתוב JSON זמני
    gen_dir = os.path.join(os.path.dirname(__file__), 'pptx_gen')
    with tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False,
                                     encoding='utf-8', dir=gen_dir) as f:
        json.dump({'title': title, 'topics': topics_data}, f, ensure_ascii=False)
        tmp_json = f.name

    tmp_pptx = tmp_json.replace('.json', '.pptx')

    try:
        result = subprocess.run(
            ['node', 'generate.js', tmp_json, tmp_pptx],
            cwd=gen_dir, capture_output=True, text=True, timeout=60
        )
        if result.returncode != 0 or not os.path.exists(tmp_pptx):
            return jsonify({'error': result.stderr or 'שגיאה ב-Node.js'}), 502

        with open(tmp_pptx, 'rb') as f:
            pptx_bytes = f.read()
    finally:
        for p in [tmp_json, tmp_pptx]:
            try: os.remove(p)
            except: pass

    safe = re.sub(r'[^\w\s-]', '', title)[:40].strip() or 'briefing'
    return send_file(
        io.BytesIO(pptx_bytes),
        as_attachment=True,
        download_name=f'{safe}.pptx',
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )




def _lang_instruction(lang):
    if lang and lang.lower().startswith('en'):
        return "Respond in English."
    return "השב בעברית."

def _is_english(lang):
    return lang and lang.lower().startswith('en')


@app.route('/api/summary', methods=['POST'])
@login_required
def ai_summary():
    data = request.json or {}
    full_text = data.get('full_text', '').strip()
    lang = data.get('lang', 'he')
    if not full_text:
        return jsonify({'error': 'אין תמלול'}), 400

    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        return jsonify({'error': 'חסר API key'}), 500

    client = anthropic.Anthropic(api_key=api_key)
    if _is_english(lang):
        prompt = f"""Read the following transcript and write a concise, structured summary in English.
The summary should include:
- 5-8 key points
- Each point as one clear sentence
- The overall topic of the episode
- Main conclusions or important ideas

Return ONLY in this format:
**Episode Topic:** [short topic]

**Key Points:**
• [point 1]
• [point 2]
• [point 3]
...

**Conclusion:** [one summary sentence]

Transcript:
{full_text[:15000]}"""
    else:
        prompt = f"""קרא את התמלול הבא וכתוב סיכום קצר ומסודר בעברית.
הסיכום צריך להיות:
- 5-8 נקודות עיקריות
- כל נקודה משפט אחד ברור
- מה הנושא הכללי של הפרק
- מה המסקנות או הרעיונות החשובים

החזר בפורמט הזה בלבד:
**נושא הפרק:** [נושא קצר]

**נקודות עיקריות:**
• [נקודה 1]
• [נקודה 2]
• [נקודה 3]
...

**מסקנה:** [משפט סיכום אחד]

התמלול:
{full_text[:15000]}"""

    try:
        msg = client.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=1000,
            messages=[{'role': 'user', 'content': prompt}]
        )
        return jsonify({'summary': msg.content[0].text.strip()})
    except Exception as e:
        return jsonify({'error': str(e)}), 502


@app.route('/api/glossary', methods=['POST'])
@login_required
def ai_glossary():
    data = request.json or {}
    full_text = data.get('full_text', '').strip()
    if not full_text:
        return jsonify({'error': 'אין תמלול'}), 400

    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        return jsonify({'error': 'חסר API key'}), 500

    lang = data.get('lang', 'he')
    client = anthropic.Anthropic(api_key=api_key)
    if _is_english(lang):
        prompt = f"""Read the following transcript and extract 5-8 key concepts/terms someone needs to know before listening to this episode.

Choose terms that are:
- Specific to the episode's topic (not general words)
- Terms that someone unfamiliar with them would struggle to follow the discussion
- Diverse — technical, conceptual, names, theories, tools, etc.

For each term, write a short and clear explanation in English (1-2 sentences), at the level of someone encountering it for the first time.

Return JSON ONLY, no extra text, in this format:
[
  {{"term": "term name", "explanation": "short clear explanation"}},
  ...
]

Transcript:
{full_text[:12000]}"""
    else:
        prompt = f"""קרא את התמלול הבא וחלץ 5-8 מושגים/מונחים מרכזיים שמישהו צריך להכיר לפני שהוא שומע את הפרק.

בחר מושגים שהם:
- ספציפיים לנושא הפרק (לא מילים כלליות)
- כאלה שמי שלא מכיר אותם יתקשה להבין את הדיון
- מגוונים — טכניים, מושגים, שמות, תיאוריות, כלים וכו'

לכל מושג כתוב הסבר קצר ובהיר בעברית (1-2 משפטים), ברמה של מישהו שנחשף אליו לראשונה.

החזר JSON בלבד, ללא טקסט נוסף, בפורמט הזה:
[
  {{"term": "שם המושג", "explanation": "הסבר קצר ובהיר"}},
  ...
]

התמלול:
{full_text[:12000]}"""

    try:
        msg = client.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=1200,
            messages=[{'role': 'user', 'content': prompt}]
        )
        raw = msg.content[0].text.strip()
        start = raw.find('[')
        end = raw.rfind(']') + 1
        terms = json.loads(raw[start:end])
        return jsonify({'terms': terms})
    except Exception as e:
        return jsonify({'error': str(e)}), 502


@app.route('/api/opportunities', methods=['POST'])
@login_required
def business_opportunities():
    data = request.json or {}
    full_text = data.get('full_text', '').strip()
    lang = data.get('lang', 'he')
    if not full_text:
        return jsonify({'error': 'אין תמלול'}), 400

    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        return jsonify({'error': 'חסר API key'}), 500

    client = anthropic.Anthropic(api_key=api_key)
    if _is_english(lang):
        prompt = f"""Read the following transcript and identify business opportunities.
For each opportunity specify:
1. Opportunity name
2. Short description
3. Target audience
4. How to implement it
5. Estimated revenue potential

Return in this format:

🚀 Opportunity 1: [name]
📌 Description: [short description]
👥 Target Audience: [who will buy]
⚙️ Implementation: [how to do it]
💰 Potential: [revenue estimate]

---

Identify 3-5 business opportunities from the content.
Transcript:
{full_text[:15000]}"""
    else:
        prompt = f"""קרא את התמלול הבא וזהה הזדמנויות עסקיות.
לכל הזדמנות ציין:
1. שם ההזדמנות
2. תיאור קצר
3. קהל יעד
4. איך לממש אותה
5. פוטנציאל הכנסה משוער

החזר בפורמט הזה:

🚀 הזדמנות 1: [שם]
📌 תיאור: [תיאור קצר]
👥 קהל יעד: [מי יקנה]
⚙️ מימוש: [איך עושים את זה]
💰 פוטנציאל: [הערכת הכנסה]

---

זהה 3-5 הזדמנויות עסקיות מהתוכן.
התמלול:
{full_text[:15000]}"""

    try:
        msg = client.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=1500,
            messages=[{'role': 'user', 'content': prompt}]
        )
        return jsonify({'opportunities': msg.content[0].text.strip()})
    except Exception as e:
        return jsonify({'error': str(e)}), 502


# ── PREMIUM PPTX ──────────────────────────────────────────────────────────────

_PALETTE = ["#7C3AED", "#F97316", "#0D9488", "#E11D48", "#4F46E5", "#059669"]

def _hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _rtl(para):
    from pptx.oxml.ns import qn
    pPr = para._p.get_or_add_pPr()
    pPr.set(qn('a:rtl'), '1')

def _rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = _hex_rgb(color)
    sh.line.fill.background()
    return sh

def _txt(slide, text, x, y, w, h, size, bold=False, color="#FFFFFF", align=PP_ALIGN.RIGHT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _rtl(p)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _hex_rgb(color)
    return tb

def _slide_title(prs, title, theme):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, 10, 7.5, "#0F172A")
    bw = 10 / len(_PALETTE)
    for i, c in enumerate(_PALETTE):
        _rect(sl, i*bw, 6.85, bw+0.02, 0.65, c)
    _txt(sl, title, 0.5, 1.5, 9, 1.6, 36, bold=True)
    _txt(sl, theme, 0.5, 3.3, 9, 1.0, 20, color="#94A3B8")
    line = sl.shapes.add_shape(1, Inches(0.5), Inches(3.1), Inches(9), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = _hex_rgb(_PALETTE[0])
    line.line.fill.background()

def _slide_content(prs, title, bullets, accent):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, 10, 7.5, "#FFFFFF")
    _rect(sl, 0, 0, 10, 1.35, "#0F172A")
    _rect(sl, 0, 0, 0.1, 7.5, accent)
    _txt(sl, title, 0.3, 0.12, 9.3, 1.1, 26, bold=True)
    for i, b in enumerate(bullets[:4]):
        y = 1.6 + i * 1.35
        dot = sl.shapes.add_shape(1, Inches(9.1), Inches(y+0.15), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = _hex_rgb(accent)
        dot.line.fill.background()
        _txt(sl, b, 0.25, y, 8.6, 1.2, 17, color="#1E293B")

def _slide_quote(prs, quote, speaker, accent):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, 10, 7.5, accent)
    _txt(sl, '❝', 0.3, 0.3, 2, 1.5, 64, bold=True, align=PP_ALIGN.LEFT)
    _txt(sl, quote, 0.6, 1.7, 8.8, 3.8, 22, bold=True)
    if speaker:
        _txt(sl, f'— {speaker}', 0.6, 5.7, 8.8, 0.7, 14, color="#E2E8F0")

def _slide_summary(prs, takeaways):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, 10, 7.5, "#0F172A")
    bw = 10 / len(_PALETTE)
    for i, c in enumerate(_PALETTE):
        _rect(sl, i*bw, 0, bw+0.02, 0.5, c)
    _txt(sl, 'תודה שהאזנת ✨', 0.5, 0.7, 9, 1, 30, bold=True)
    for i, t in enumerate(takeaways[:3]):
        y = 2.0 + i * 1.55
        _rect(sl, 0.4, y, 9.2, 1.25, _PALETTE[i])
        _txt(sl, t, 0.65, y+0.18, 8.7, 0.9, 15)

def _claude_extract(full_text, title):
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        return None
    client = anthropic.Anthropic(api_key=api_key)
    prompt = f"""אתה עוזר שיוצר מצגות מקצועיות מתמלולי פודקאסטים.
קרא את התמלול הבא וצור מבנה JSON למצגת.

חוקים חשובים:
- 10-14 slides בסה"כ
- לפחות 2-3 slides מסוג "quote" (ציטוטים מרשימים מהפרק)
- כותרת: עד 6 מילים, ברורה
- כל bullet: משפט שלם 15-25 מילים
- 3-4 bullets לכל content slide
- episode_theme: משפט אחד המסכם את כל הפרק
- takeaways: 3 לקחים מרכזיים

החזר JSON בלבד:
{{
  "episode_theme": "...",
  "takeaways": ["...", "...", "..."],
  "slides": [
    {{"type": "content", "title": "...", "bullets": ["...", "...", "..."]}},
    {{"type": "quote", "quote": "...", "speaker": "..."}}
  ]
}}

כותרת: {title}
תמלול:
{full_text[:14000]}"""

    msg = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=3000,
        messages=[{'role': 'user', 'content': prompt}]
    )
    raw = msg.content[0].text.strip()
    m = re.search(r'\{[\s\S]*\}', raw)
    if m:
        return json.loads(m.group())
    return None


@app.route('/api/premium_pptx', methods=['POST'])
@login_required
def premium_pptx():
    data = request.json or {}
    full_text = data.get('full_text', '').strip()
    title = data.get('title', 'סיכום פרק')
    if not full_text:
        return jsonify({'error': 'אין תמלול'}), 400

    slides_data = _claude_extract(full_text, title)
    if not slides_data:
        return jsonify({'error': 'שגיאה בחילוץ תוכן'}), 502

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _slide_title(prs, title, slides_data.get('episode_theme', ''))

    for idx, s in enumerate(slides_data.get('slides', [])):
        accent = _PALETTE[idx % len(_PALETTE)]
        if s.get('type') == 'quote':
            _slide_quote(prs, s.get('quote', ''), s.get('speaker', ''), accent)
        else:
            _slide_content(prs, s.get('title', ''), s.get('bullets', []), accent)

    _slide_summary(prs, slides_data.get('takeaways', []))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    safe = re.sub(r'[^\w\s-]', '', title)[:40].strip()
    return send_file(buf, as_attachment=True,
                     download_name=f'{safe}.pptx',
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')


@app.route('/api/create_post', methods=['POST'])
@login_required
def create_post():
    data = request.json or {}
    full_text = data.get('full_text', '').strip()
    title = data.get('title', '')
    if not full_text:
        return jsonify({'error': 'אין תמלול'}), 400
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        return jsonify({'error': 'חסר API key'}), 500
    client = anthropic.Anthropic(api_key=api_key)
    prompt = f"""כתוב פוסט מקצועי לרשתות חברתיות (לינקדאין / פייסבוק) בעברית, מבוסס על הפרק הבא.

הפוסט צריך להיות:
- פותח בשורה מושכת שגורמת לאנשים לעצור ולקרוא
- 3-4 תובנות מרכזיות מהפרק (כל אחת בשורה נפרדת עם אימוג׳י)
- משפט סיום שמזמין אנשים לדון / לשתף / לשאול
- האשטגים רלוונטיים בסוף (4-6)
- אורך: 150-250 מילים

כותרת הפרק: {title}
תמלול:
{full_text[:10000]}

כתוב רק את הפוסט עצמו, ללא הסברים."""

    msg = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=1000,
        messages=[{'role': 'user', 'content': prompt}]
    )
    return jsonify({'post': msg.content[0].text.strip()})


@app.route('/api/ask', methods=['POST'])
@login_required
def ask_episode():
    data = request.json or {}
    full_text = data.get('full_text', '').strip()
    question = data.get('question', '').strip()
    lang = data.get('lang', 'he')
    if not full_text or not question:
        return jsonify({'error': 'חסר תמלול או שאלה'}), 400
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        return jsonify({'error': 'חסר API key'}), 500
    client = anthropic.Anthropic(api_key=api_key)
    if _is_english(lang):
        content = f"""Answer the following question in English, based only on the transcript.
If the answer is not in the transcript — say so clearly.
Answer briefly and to the point (2-4 sentences).

Question: {question}

Transcript:
{full_text[:12000]}"""
    else:
        content = f"""ענה על השאלה הבאה בעברית, בהתבסס על התמלול בלבד.
אם התשובה לא נמצאת בתמלול — אמור זאת בבירור.
ענה בקצרה ולעניין (2-4 משפטים).

שאלה: {question}

תמלול:
{full_text[:12000]}"""
    msg = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=600,
        messages=[{'role': 'user', 'content': content}]
    )
    return jsonify({'answer': msg.content[0].text.strip()})


@app.route('/api/share', methods=['POST'])
@login_required
def create_share():
    data = request.json or {}
    video_id     = data.get('video_id', '').strip()
    title        = data.get('title', '').strip()
    summary      = data.get('summary', '').strip()
    opportunities = data.get('opportunities', '').strip()
    lang         = data.get('lang', 'he')
    if not video_id:
        return jsonify({'error': 'חסר video_id'}), 400
    share_id = uuid4().hex[:12]
    user_db.create_share(share_id, video_id, title, summary, opportunities, lang)
    return jsonify({'share_id': share_id, 'url': f'/share/{share_id}'})


@app.route('/share/<share_id>')
def view_share(share_id):
    share = user_db.get_share(share_id)
    if not share:
        return render_template('landing.html'), 404
    return render_template('share.html', share=share)


@app.route('/en')
def landing_en():
    return render_template('landing_en.html')


@app.route('/api/save_episode', methods=['POST'])
@login_required
def save_episode():
    u = current_user()
    data = request.json or {}
    title = data.get('title', '').strip()
    video_id = data.get('video_id', '').strip()
    summary = data.get('summary', '').strip()
    if not title or not video_id:
        return jsonify({'error': 'חסר מידע'}), 400
    user_db.save_episode(u['id'], video_id, title, summary)
    return jsonify({'success': True})


@app.route('/api/history', methods=['GET'])
@login_required
def get_history():
    u = current_user()
    episodes = user_db.get_episodes(u['id'])
    return jsonify({'episodes': episodes})


@app.route('/api/upload_cookies', methods=['POST'])
@login_required
def upload_cookies():
    """העלאת קובץ cookies של YouTube — מאפשר גישה מהשרת הענן"""
    if 'file' not in request.files:
        return jsonify({'error': 'לא נבחר קובץ'}), 400
    f = request.files['file']
    if not f.filename.endswith('.txt'):
        return jsonify({'error': 'הקובץ חייב להיות .txt'}), 400
    f.save(COOKIES_FILE)
    return jsonify({'success': True, 'message': 'Cookies הועלו בהצלחה!'})


@app.route('/api/send_email', methods=['POST'])
@login_required
def send_email():
    """שלח סיכום למייל דרך Resend"""
    data = request.json or {}
    to_email = data.get('email', '').strip()
    summary  = data.get('summary', '').strip()
    title    = data.get('title', 'סיכום פרק').strip()

    if not to_email or not summary:
        return jsonify({'error': 'חסר מייל או סיכום'}), 400

    if not RESEND_KEY:
        return jsonify({'error': 'מפתח Resend לא מוגדר'}), 500

    # המר סיכום ל-HTML פשוט
    html_lines = []
    for line in summary.split('\n'):
        line = line.strip()
        if not line:
            html_lines.append('<br>')
        elif line.startswith('**') and line.endswith('**'):
            html_lines.append(f'<h3 style="color:#1a1a2e">{line[2:-2]}</h3>')
        elif line.startswith('•'):
            html_lines.append(f'<li style="margin:6px 0">{line[1:].strip()}</li>')
        else:
            html_lines.append(f'<p>{line}</p>')

    html_body = f"""
    <div dir="rtl" style="font-family:Arial,sans-serif;max-width:600px;margin:0 auto;padding:24px">
      <div style="background:linear-gradient(135deg,#1a0533,#2d1060);padding:20px;border-radius:10px;margin-bottom:20px">
        <h1 style="color:#fff;margin:0;font-size:1.4rem">שולחן ארבע</h1>
        <p style="color:#a78bfa;margin:4px 0 0">סיכום אוטומטי</p>
      </div>
      <h2 style="color:#1a0533">{title}</h2>
      {''.join(html_lines)}
      <hr style="border:none;border-top:1px solid #eee;margin:24px 0">
      <p style="color:#999;font-size:0.75rem">נוצר על ידי שולחן ארבע · AI Summary</p>
    </div>"""

    try:
        resp = http_requests.post(
            'https://api.resend.com/emails',
            headers={
                'Authorization': f'Bearer {RESEND_KEY}',
                'Content-Type': 'application/json'
            },
            json={
                'from': 'שולחן ארבע <onboarding@resend.dev>',
                'to': [to_email],
                'subject': f'סיכום: {title}',
                'html': html_body
            },
            timeout=15
        )
        if resp.status_code == 200:
            return jsonify({'success': True})
        return jsonify({'error': resp.text}), resp.status_code
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/telegram-webhook', methods=['POST'])
def telegram_webhook():
    """מקבל callback מ-Telegram כשמנכ"ל לוחץ על כפתור אישור."""
    import ceo_bot
    payload = request.json or {}

    def _refund(uid):
        notify_ceo(f'💸 החזר כספי בוצע למשתמש {uid}')

    actions = {
        f'refund:{uid}': (lambda u=uid: _refund(u))
        for uid in [payload.get('callback_query', {}).get('data', '').split(':')[-1]]
        if uid.isdigit()
    }

    ceo_bot.handle_webhook_payload(payload, actions)
    return '', 200


LI_CLIENT_ID         = os.getenv('LINKEDIN_CLIENT_ID', '')
LI_CLIENT_SECRET     = os.getenv('LINKEDIN_CLIENT_SECRET', '')
LI_REDIRECT_URI      = os.getenv('APP_URL', '') + '/linkedin/callback'
LI_SCOPES            = 'openid profile w_member_social'
LINKEDIN_ADMIN_TOKEN = os.getenv('LINKEDIN_ADMIN_TOKEN', '')
LINKEDIN_ADMIN_URN   = os.getenv('LINKEDIN_ADMIN_URN', '')

LINKEDIN_STYLES = {
    'A': 'התחל את הפוסט עם שאלה פרובוקטיבית שגורמת לקורא לעצור ולחשוב',
    'B': 'התחל עם נתון מפתיע או סטטיסטיקה בלתי צפויה מהתוכן',
    'C': 'התחל עם מיקרו-סיפור קצר (2 משפטים) שמחבר את הנושא לחיי היומיום',
}
LINKEDIN_STYLES_EN = {
    'A': 'Start with a provocative question that makes the reader stop and think',
    'B': 'Start with a surprising stat or unexpected data point from the content',
    'C': 'Start with a micro-story (2 sentences) that connects the topic to everyday life',
}


@app.route('/li-debug')
def li_debug():
    return jsonify({
        'client_id_set': bool(LI_CLIENT_ID),
        'client_id_len': len(LI_CLIENT_ID),
        'secret_set': bool(LI_CLIENT_SECRET),
        'secret_len': len(LI_CLIENT_SECRET),
        'redirect_uri': LI_REDIRECT_URI,
        'scopes': LI_SCOPES,
    })

@app.route('/linkedin/connect')
@login_required
def linkedin_connect():
    u = current_user()
    if not is_admin(u):
        return redirect('/')
    import urllib.parse
    params = urllib.parse.urlencode({
        'response_type': 'code',
        'client_id': LI_CLIENT_ID,
        'redirect_uri': LI_REDIRECT_URI,
        'scope': LI_SCOPES,
        'state': session['user_id'],
    })
    return redirect(f'https://www.linkedin.com/oauth/v2/authorization?{params}')


@app.route('/linkedin/callback')
def linkedin_callback():
    # Use state param (= user_id sent in connect) in case session cookie is lost during OAuth redirect
    state_uid = request.args.get('state')
    u = user_db.get_by_id(state_uid) if state_uid else current_user()
    print(f'[linkedin/callback] state_uid={state_uid} user={u and u.get("email")} session_uid={session.get("user_id")}')
    if not is_admin(u):
        print(f'[linkedin/callback] not admin, redirecting')
        return redirect('/')
    code = request.args.get('code')
    error = request.args.get('error')
    error_desc = request.args.get('error_description', '')
    print(f'[linkedin/callback] code={bool(code)} error={error} desc={error_desc}')
    if error or not code:
        import urllib.parse as _up
        return redirect(f'/?linkedin_error=1&detail={_up.quote(error_desc or error or "no_code")}')

    # החלף code ב-access token
    try:
        token_resp = http_requests.post(
            'https://www.linkedin.com/oauth/v2/accessToken',
            data={
                'grant_type': 'authorization_code',
                'code': code,
                'redirect_uri': LI_REDIRECT_URI,
                'client_id': LI_CLIENT_ID,
                'client_secret': LI_CLIENT_SECRET,
            },
            headers={'Content-Type': 'application/x-www-form-urlencoded'},
            timeout=10
        )
        token_resp.raise_for_status()
        token_data = token_resp.json()
        access_token = token_data['access_token']
        expires_in = token_data.get('expires_in', 5184000)

        # קבל URN של המשתמש
        me_resp = http_requests.get(
            'https://api.linkedin.com/v2/userinfo',
            headers={'Authorization': f'Bearer {access_token}'},
            timeout=10
        )
        me_resp.raise_for_status()
        urn = me_resp.json().get('sub', '')

        from datetime import timedelta
        expiry = datetime.now() + timedelta(seconds=expires_in)
        uid = u['id'] if u else session.get('user_id')
        if uid:
            user_db.save_linkedin_token(uid, access_token, expiry, f'urn:li:person:{urn}')
            print(f'[linkedin/callback] saved token for uid={uid} urn={urn}')
    except Exception as e:
        print(f'[linkedin/callback] {e}')
        return redirect('/?linkedin_error=1')

    return redirect('/?linkedin_connected=1')


@app.route('/api/linkedin/post', methods=['POST'])
@login_required
@limiter.limit("5/minute", key_func=lambda: str(session.get('user_id', get_remote_address())))
def linkedin_post():
    u = current_user()
    if not is_admin(u):
        return jsonify({'error': 'Unauthorized'}), 403
    # Use admin token if set, otherwise fall back to per-user OAuth token
    li_token = LINKEDIN_ADMIN_TOKEN or u.get('linkedin_access_token', '')
    li_urn   = LINKEDIN_ADMIN_URN   or u.get('linkedin_urn', '')
    if not li_token:
        return jsonify({'error': 'לא מחובר ל-LinkedIn', 'connect': True}), 403

    data = request.json or {}
    full_text = data.get('full_text', '').strip()
    video_id  = data.get('video_id', '').strip()
    lang      = data.get('lang', 'he')
    if not full_text:
        return jsonify({'error': 'אין תוכן'}), 400

    # בחר סגנון לפי היסטוריה (round-robin בשלב הראשון)
    stats = user_db.get_linkedin_stats(u['id'])
    used  = {s['style']: s['posts'] for s in stats}
    style = min(['A', 'B', 'C'], key=lambda s: used.get(s, 0))
    styles_dict = LINKEDIN_STYLES_EN if _is_english(lang) else LINKEDIN_STYLES
    style_instruction = styles_dict[style]

    api_key = os.getenv('ANTHROPIC_API_KEY')
    ai = anthropic.Anthropic(api_key=api_key)
    if _is_english(lang):
        prompt = f"""Write a LinkedIn post in English using an interesting insight from this episode as the opening hook, then promote the app getpodsnap.com.

Style instruction: {style_instruction}

Post structure:
1. Strong opening — a surprising insight, stat, or quote from the episode (1-2 sentences)
2. Brief expansion on the topic (2-3 bullet points •)
3. Natural transition: "I analyzed this episode with AI..." or "I ran it through PodSnap and..."
4. Promote getpodsnap.com — an app that analyzes podcasts with AI and produces summaries, business opportunities, glossary
5. CTA: "Try it free" / "Link in comments" / question to audience
6. 3-5 relevant hashtags

Requirements:
- Length: 150-250 words
- The post should feel authentic, not like a direct ad
- Write the post directly, no title or explanation

Episode content:
{full_text[:8000]}"""
        footer = "\n\n🎙️ Full AI podcast analysis → getpodsnap.com"
    else:
        prompt = f"""כתוב פוסט LinkedIn בעברית שמשתמש בתובנה מעניינת מהפרק כ-hook פותח, ואז מפרסם את האפליקציה getpodsnap.com.

הנחיית סגנון: {style_instruction}

מבנה הפוסט:
1. פתיחה חזקה — תובנה, נתון, או ציטוט מפתיע מהפרק (1-2 משפטים)
2. הרחבה קצרה על הנושא (2-3 bullet points •)
3. מעבר טבעי: "ניתחתי את הפרק הזה עם AI..." או "הכנסתי את זה ל-PodSnap ו..."
4. פרסום של getpodsnap.com — אפליקציה שמנתחת פודקאסטים עם AI ומוציאה סיכום, הזדמנויות עסקיות, מילון מושגים
5. CTA: "תנסו בחינם" / "קישור בתגובות" / שאלה לקהל
6. 3-5 hashtags

דרישות:
- אורך: 150-250 מילים
- הפוסט צריך להרגיש אותנטי, לא כמו פרסומת ישירה
- כתוב ישירות את הפוסט, ללא כותרת או הסבר

תוכן הפרק:
{full_text[:8000]}"""
        footer = "\n\n🎙️ ניתוח AI מלא של פודקאסטים → getpodsnap.com"

    try:
        msg = ai.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=600,
            messages=[{'role': 'user', 'content': prompt}]
        )
        post_text = msg.content[0].text.strip()
        post_text += footer
    except Exception as e:
        return jsonify({'error': f'שגיאת AI: {str(e)}'}), 502

    # פרסם ל-LinkedIn
    try:
        li_resp = http_requests.post(
            'https://api.linkedin.com/v2/ugcPosts',
            headers={
                'Authorization': f'Bearer {li_token}',
                'Content-Type': 'application/json',
                'X-Restli-Protocol-Version': '2.0.0',
            },
            json={
                'author': li_urn,
                'lifecycleState': 'PUBLISHED',
                'specificContent': {
                    'com.linkedin.ugc.ShareContent': {
                        'shareCommentary': {'text': post_text},
                        'shareMediaCategory': 'NONE',
                    }
                },
                'visibility': {
                    'com.linkedin.ugc.MemberNetworkVisibility': 'PUBLIC'
                }
            },
            timeout=15
        )
        li_resp.raise_for_status()
        post_urn = li_resp.json().get('id', '')
        user_db.save_linkedin_post(u['id'], post_urn, post_text, video_id, style)
        post_url = f'https://www.linkedin.com/feed/update/{post_urn}/' if post_urn else ''
        return jsonify({'success': True, 'post': post_text, 'url': post_url, 'style': style})
    except Exception as e:
        return jsonify({'error': f'שגיאת LinkedIn: {str(e)}'}), 502


@app.route('/api/admin/run-agent', methods=['POST'])
@login_required
def run_agent():
    u = current_user()
    if not is_admin(u):
        return jsonify({'error': 'Unauthorized'}), 403
    try:
        import run_linkedin_agent
        import importlib
        importlib.reload(run_linkedin_agent)
        import io, contextlib
        buf = io.StringIO()
        with contextlib.redirect_stdout(buf):
            run_linkedin_agent.run()
        return jsonify({'success': True, 'log': buf.getvalue()})
    except Exception as e:
        import traceback
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


@app.route('/api/admin/agent-posts')
@login_required
def agent_posts():
    u = current_user()
    if not is_admin(u):
        return jsonify({'error': 'Unauthorized'}), 403
    posts = user_db.get_auto_posts(10)
    return jsonify({'posts': [dict(p, posted_at=str(p['posted_at'])) for p in posts]})


@app.errorhandler(500)
def server_error(e):
    return jsonify({'error': f'Server error: {str(e)}'}), 500


@app.errorhandler(Exception)
def unhandled(e):
    import traceback
    print(traceback.format_exc())
    return jsonify({'error': str(e)}), 500


@app.errorhandler(429)
def rate_limit_handler(e):
    return jsonify({'error': 'יותר מדי בקשות — נסה שוב בעוד דקה'}), 429


@app.route('/health')
def health():
    try:
        user_db.ping()
        return jsonify({'status': 'ok'})
    except Exception as e:
        return jsonify({'status': 'db_down', 'error': str(e)}), 500


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5002))
    print(f"SikumAI -> http://localhost:{port}")
    app.run(debug=False, host='0.0.0.0', port=port)
